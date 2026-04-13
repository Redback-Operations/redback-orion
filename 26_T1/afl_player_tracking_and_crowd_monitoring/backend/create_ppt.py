from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
RED = RGBColor(0xC0, 0x00, 0x00)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xFF, 0x80, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=LIGHT_GRAY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def textbox(slide, text, l, t, w, h, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.4, DARK_RED)
    textbox(slide, title, 0.3, 0.1, 10, 0.8, size=32, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, subtitle, 0.3, 0.85, 10, 0.5, size=14, color=RGBColor(0xFF, 0xCC, 0xCC))


def bullet_box(slide, title, bullets, l, t, w, h, title_color=DARK_RED, bullet_size=15):
    box(slide, l, t, w, h, WHITE)
    textbox(slide, title, l + 0.15, t + 0.1, w - 0.3, 0.4, size=16, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.55), Inches(w - 0.3), Inches(h - 0.65))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = DARK_GRAY


# ─── SLIDE 1: TITLE ───────────────────────────────────────────────
slide = add_slide()
bg(slide, DARK_RED)
box(slide, 0, 2.5, 13.33, 0.08, WHITE)
textbox(slide, "Project Orion", 1, 0.6, 11, 1.2, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "AFL Player Tracking & Crowd Monitoring", 1, 1.7, 11, 0.7, size=24, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
textbox(slide, "Backend Team — Sprint 1 Progress Report", 1, 2.7, 11, 0.6, size=20, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "SIT782 – Team Project (B) – Execution and Delivery", 1, 3.4, 11, 0.5, size=14, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
textbox(slide, "Trimester 1, 2026  |  Deakin University", 1, 5.8, 11, 0.5, size=13, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)


# ─── SLIDE 2: TEAM ────────────────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "Backend Team", "Who we are and what we own")

members = [
    ("Tomin Jose", "Backend Lead", "Architecture, CORS, route wiring, auth core"),
    ("Lucas Targett", "Auth Developer", "API routes, auth schemas, auth endpoints"),
    ("William Hamilton", "Service Layer", "Mock services, player & crowd clients"),
    ("Prabhnoor", "DB & Infra", "Config setup, database models, error handling"),
]

for i, (name, role, task) in enumerate(members):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.6 + row * 2.5
    box(slide, x, y, 6.2, 2.2, WHITE)
    box(slide, x, y, 6.2, 0.45, DARK_RED)
    textbox(slide, name, x + 0.15, y + 0.05, 5.5, 0.38, size=15, bold=True, color=WHITE)
    textbox(slide, role, x + 0.15, y + 0.55, 5.5, 0.35, size=13, bold=True, color=DARK_RED)
    textbox(slide, task, x + 0.15, y + 0.95, 5.8, 0.9, size=12, color=DARK_GRAY)


# ─── SLIDE 3: ARCHITECTURE ────────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "System Architecture", "How the backend fits into Project Orion")

# Frontend
box(slide, 0.5, 1.6, 2.5, 0.8, RGBColor(0x20, 0x60, 0xA0))
textbox(slide, "Frontend (React)", 0.5, 1.6, 2.5, 0.8, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "localhost:3000", 0.5, 2.3, 2.5, 0.4, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Arrow down
textbox(slide, "▼", 1.45, 2.7, 0.8, 0.4, size=18, color=DARK_RED, align=PP_ALIGN.CENTER)

# Backend Gateway
box(slide, 0.3, 3.1, 2.9, 1.0, DARK_RED)
textbox(slide, "Backend API Gateway", 0.3, 3.15, 2.9, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "FastAPI  |  localhost:8000", 0.3, 3.6, 2.9, 0.4, size=11, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
textbox(slide, "OUR RESPONSIBILITY", 0.3, 4.1, 2.9, 0.35, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows
textbox(slide, "▼                    ▼", 2.2, 4.2, 3, 0.4, size=16, color=DARK_RED)

# Player Service
box(slide, 3.4, 4.7, 2.5, 0.8, RGBColor(0x20, 0x80, 0x40))
textbox(slide, "Player Service (YOLO)", 3.4, 4.7, 2.5, 0.5, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "localhost:8001", 3.4, 5.15, 2.5, 0.35, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Crowd Service
box(slide, 6.2, 4.7, 2.5, 0.8, RGBColor(0x80, 0x40, 0xA0))
textbox(slide, "Crowd Service (Density)", 6.2, 4.7, 2.5, 0.5, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "localhost:8002", 6.2, 5.15, 2.5, 0.35, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# DB
box(slide, 9.2, 3.1, 2.5, 1.0, RGBColor(0x20, 0x60, 0xA0))
textbox(slide, "PostgreSQL", 9.2, 3.2, 2.5, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "Users + Jobs", 9.2, 3.65, 2.5, 0.35, size=11, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Key points
bullet_box(slide, "Key Design Decisions", [
    "asyncio.gather() — parallel calls to player & crowd",
    "USE_MOCK_SERVICES flag — toggle mocks without code changes",
    "BackgroundTasks — non-blocking video processing",
    "JWT + bcrypt — secure authentication",
    "Docker in Phase 4 only — local dev first",
], 9.2, 1.6, 3.9, 3.2)


# ─── SLIDE 4: API ENDPOINTS ───────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "API Endpoints", "What the backend exposes")

endpoints = [
    ("Auth (Public)", [
        "POST  /auth/register  — create user account",
        "POST  /auth/login     — returns JWT token",
        "GET   /auth/me        — current logged-in user",
    ], 0.3, 1.6, 4.1, 2.2),
    ("Video Processing (Protected)", [
        "POST  /upload              — accepts video, returns job_id",
        "GET   /status/{job_id}    — poll for job progress",
        "GET   /jobs               — paginated past jobs",
        "GET   /jobs/{job_id}      — full job details",
        "POST  /jobs/{job_id}/retry — retry failed service",
        "DELETE /jobs/{job_id}     — delete a job",
    ], 0.3, 3.9, 4.1, 3.0),
    ("System", [
        "GET  /health  — checks gateway + both services",
        "GET  /         — root health check",
        "GET  /api/players — mock player data",
        "GET  /api/crowd   — mock crowd data",
    ], 4.6, 1.6, 4.1, 2.2),
    ("Response Format", [
        "status: processing / done / partial / failed",
        "Partial results saved if one service fails",
        "Retry endpoint re-runs only the failed service",
        "All protected routes require JWT token",
    ], 4.6, 3.9, 4.1, 2.2),
    ("Tech Stack", [
        "FastAPI + Python 3.13",
        "SQLAlchemy + PostgreSQL",
        "httpx — async HTTP client",
        "bcrypt + python-jose (JWT)",
        "pytest + httpx (testing)",
    ], 9.0, 1.6, 4.0, 2.5),
]

for title, bullets, l, t, w, h in endpoints:
    bullet_box(slide, title, bullets, l, t, w, h)


# ─── SLIDE 5: WEEK 1 PROGRESS ─────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "Sprint 1 — Week 1 Progress", "23rd – 28th March 2026")

tasks = [
    ("Tomin", "Backend Architecture & Setup",
     ["FastAPI project structure created", "main.py with base routes (/, /health)", "BACKEND_PLAN.txt shared with team", "GIT_GUIDE.txt created for team"],
     GREEN, 0.3, 1.6),
    ("Lucas", "API Routes Implementation",
     ["routes/players.py — /api/players", "routes/crowd.py — /api/crowd", "routes/test.py — /api/test", "All routes registered in main.py"],
     GREEN, 6.5, 1.6),
    ("William", "Mock Service Layer",
     ["player_client.py with mock + real toggle", "crowd_client.py with mock + real toggle", "USE_MOCK_SERVICES flag wired", "Routes call clients (not hardcoded)"],
     GREEN, 0.3, 4.0),
    ("Prabhnoor", "Configuration Setup",
     ["config.py — loads all env vars", ".env — all local dev values set", "USE_MOCK_SERVICES=true as default", "python-dotenv integrated"],
     GREEN, 6.5, 4.0),
]

for member, task, bullets, status_color, x, y in tasks:
    box(slide, x, y, 6.5, 2.9, WHITE)
    box(slide, x, y, 6.5, 0.45, DARK_RED)
    textbox(slide, f"{member} — {task}", x + 0.15, y + 0.05, 6.1, 0.38, size=13, bold=True, color=WHITE)
    box(slide, x + 5.7, y + 0.55, 0.6, 0.35, status_color)
    textbox(slide, "DONE", x + 5.7, y + 0.55, 0.6, 0.35, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.6), Inches(5.5), Inches(2.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY


# ─── SLIDE 6: WEEK 2 PROGRESS ─────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "Sprint 1 — Week 2 Progress", "29th March – 3rd April 2026")

tasks2 = [
    ("Tomin", "CORS + Route Wiring",
     ["CORS configured for localhost:3000", "Auth router wired into main.py at /auth", "All routes registered and working", "Verified with live server run"],
     GREEN, "DONE", 0.3, 1.6),
    ("Lucas", "Auth Schemas",
     ["schemas/auth.py — RegisterRequest, LoginRequest", "schemas/auth.py — AuthResponse, UserResponse", "schemas/health.py — HealthResponse", "Merged via backend-feature/auth-schemas branch"],
     GREEN, "DONE", 6.5, 1.6),
    ("William", "Mock Services Connected to Routes",
     ["player_client + crowd_client wired to routes", "Routes no longer hardcode data", "USE_MOCK_SERVICES=true returns mock JSON", "Verified endpoints returning correct data"],
     GREEN, "DONE", 0.3, 4.0),
    ("Prabhnoor", "Database Setup",
     ["database.py — SQLAlchemy engine + session", "models.py — User + Job table definitions", "PostgreSQL running on localhost:5432", "users + jobs tables created successfully"],
     GREEN, "DONE", 6.5, 4.0),
]

for member, task, bullets, status_color, status_text, x, y in tasks2:
    box(slide, x, y, 6.5, 2.9, WHITE)
    box(slide, x, y, 6.5, 0.45, DARK_RED)
    textbox(slide, f"{member} — {task}", x + 0.15, y + 0.05, 6.1, 0.38, size=13, bold=True, color=WHITE)
    box(slide, x + 5.1, y + 0.55, 1.2, 0.35, status_color)
    textbox(slide, status_text, x + 5.1, y + 0.55, 1.2, 0.35, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.6), Inches(5.5), Inches(2.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{'✓' if status_text == 'DONE' else '⏳'}  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY


# ─── SLIDE 7: LIVE DEMO / WHAT'S WORKING ──────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "What's Working Right Now", "Live demo — backend running on localhost:8000")

working = [
    ("GET  /", '{ "status": "success", "message": "Backend is running!" }'),
    ("GET  /health", '{ "gateway": "ok", "player_service": "pending", "crowd_service": "pending" }'),
    ("GET  /api/players", '{ "players": [ { "player_id": 1, "team": "Team A", "speed": 6.4, ... } ] }'),
    ("GET  /api/crowd", '{ "crowd_count": 15432, "density": 0.72, "sections": [...] }'),
    ("POST /auth/register", '{ "access_token": "mock_token", "token_type": "bearer" }'),
    ("POST /auth/login", '{ "access_token": "mock_token", "token_type": "bearer" }'),
]

for i, (endpoint, response) in enumerate(working):
    row = i % 3
    col = i // 3
    x = 0.3 + col * 6.5
    y = 1.6 + row * 1.85
    box(slide, x, y, 6.3, 1.7, WHITE)
    box(slide, x, y, 6.3, 0.4, DARK_RED)
    textbox(slide, endpoint, x + 0.15, y + 0.05, 6.0, 0.35, size=13, bold=True, color=WHITE)
    textbox(slide, response, x + 0.15, y + 0.5, 6.0, 1.1, size=11, color=DARK_GRAY)

textbox(slide, "Swagger UI available at: http://localhost:8000/docs", 0.3, 7.1, 12, 0.35, size=12, bold=True, color=DARK_RED, align=PP_ALIGN.CENTER)


# ─── SLIDE 8: BRANCH STRUCTURE ────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "Git Branch Structure", "How the team manages code")

branches = [
    ("backend-api-gateway", "Tomin", "Main backend branch — all reviewed code merged here", DARK_RED),
    ("backend-feature/auth-schemas", "Lucas", "Auth schemas + health schema (merged)", RGBColor(0x20, 0x60, 0xA0)),
    ("feature/mock-services", "William", "Mock player & crowd services (merged via PR #2)", RGBColor(0x20, 0x80, 0x40)),
    ("config-error-handling", "Prabhnoor", "Config + error handling (not yet pushed to remote)", ORANGE),
]

for i, (branch, owner, desc, color) in enumerate(branches):
    y = 1.7 + i * 1.35
    box(slide, 0.3, y, 12.7, 1.2, WHITE)
    box(slide, 0.3, y, 0.3, 1.2, color)
    textbox(slide, branch, 0.75, y + 0.05, 5.5, 0.45, size=14, bold=True, color=color)
    textbox(slide, f"Owner: {owner}", 0.75, y + 0.5, 3.0, 0.35, size=12, color=DARK_GRAY)
    textbox(slide, desc, 5.5, y + 0.25, 7.3, 0.7, size=12, color=DARK_GRAY)

textbox(slide, "Workflow: Each member works on their branch → raises PR → Tomin reviews → merges into backend-api-gateway", 0.3, 7.0, 12.7, 0.4, size=12, bold=True, color=DARK_RED, align=PP_ALIGN.CENTER)


# ─── SLIDE 9: NEXT STEPS ──────────────────────────────────────────
slide = add_slide()
bg(slide)
header_bar(slide, "Week 3 Progress", "Sprint 2: 4th – 10th April 2026")

next_steps = [
    ("Tomin — Week 3", [
        "auth/hashing.py — bcrypt hash + verify",
        "auth/jwt.py — create + decode JWT tokens",
        "auth/dependencies.py — get_current_user + require_admin",
        "Contract template sent to player & crowd teams",
    ], 0.3, 1.6),
    ("Lucas — Week 3", [
        "POST /auth/register — real DB implementation",
        "POST /auth/login — verify password + issue JWT",
        "GET  /auth/me — protected endpoint",
    ], 6.7, 1.6),
    ("William — Week 3", [
        "schemas/jobs.py — JobSummary",
        "schemas/jobs.py — JobDetail",
        "schemas/jobs.py — JobListResponse",
        "schemas/jobs.py — UploadResponse",
    ], 0.3, 4.0),
    ("Prabhnoor — Week 3", [
        "Global exception handler in main.py",
        "Request logging middleware",
        "All unhandled errors return clean JSON",
    ], 6.7, 4.0),
]

status_map = {
    "Tomin — Week 3": (GREEN, "DONE"),
    "Lucas — Week 3": (ORANGE, "IN PROGRESS"),
    "William — Week 3": (ORANGE, "IN PROGRESS"),
    "Prabhnoor — Week 3": (ORANGE, "IN PROGRESS"),
}

for title, bullets, x, y in next_steps:
    status_color, status_text = status_map[title]
    box(slide, x, y, 6.2, 2.9, WHITE)
    box(slide, x, y, 6.2, 0.45, DARK_RED)
    textbox(slide, title, x + 0.15, y + 0.05, 5.5, 0.38, size=13, bold=True, color=WHITE)
    box(slide, x + 4.8, y + 0.55, 1.2, 0.35, status_color)
    textbox(slide, status_text, x + 4.8, y + 0.55, 1.2, 0.35, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.6), Inches(5.5), Inches(2.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{'Done' if status_text == 'DONE' else 'TODO'}  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY


# ─── SLIDE 10: THANK YOU ──────────────────────────────────────────
slide = add_slide()
bg(slide, DARK_RED)
textbox(slide, "Thank You", 1, 1.5, 11, 1.5, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "Project Orion — Backend Team", 1, 3.0, 11, 0.7, size=22, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
textbox(slide, "Questions?", 1, 3.8, 11, 0.6, size=20, color=WHITE, align=PP_ALIGN.CENTER)
box(slide, 3, 4.6, 7.33, 0.06, WHITE)
textbox(slide, "Tomin  |  Lucas  |  William  |  Prabhnoor", 1, 4.8, 11, 0.5, size=15, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
textbox(slide, "SIT782 – Deakin University – Trimester 1, 2026", 1, 5.4, 11, 0.5, size=13, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)


output = "d:/Masters/trimester 1 2026/SIT782 – Team Project (B) – Execution and Delivery/redback-orion/26_T1/afl_player_tracking_and_crowd_monitoring/backend/Project_Orion_Backend_Progress.pptx"
prs.save(output)
print(f"Saved: {output}")
