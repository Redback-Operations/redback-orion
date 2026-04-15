from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x2B, 0x55)   # primary dark
MID_NAVY  = RGBColor(0x1A, 0x48, 0x8A)   # accent blue
TEAL      = RGBColor(0x00, 0x8C, 0x99)   # highlight
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # accent yellow
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF9, 0xFC)   # slide background
LIGHT_BLUE= RGBColor(0xE3, 0xEE, 0xF9)   # card background
MID_GRAY  = RGBColor(0x88, 0x99, 0xAA)   # secondary text
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)   # body text
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED       = RGBColor(0xC0, 0x39, 0x2B)
DARK_RED  = NAVY
BLUE      = MID_NAVY
BLACK     = DARK_GRAY

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=OFF_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def textbox(slide, text, l, t, w, h, size=16, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.4, NAVY)
    box(slide, 0, 1.35, 13.33, 0.08, TEAL)   # teal accent line
    textbox(slide, title, 0.35, 0.08, 12, 0.78, size=30, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, subtitle, 0.35, 0.85, 12.5, 0.42,
                size=13, color=RGBColor(0xAA, 0xCC, 0xEE))


def card(slide, title, bullets, l, t, w, h,
         title_color=NAVY, bullet_size=14, title_size=15):
    box(slide, l, t, w, h, LIGHT_BLUE)
    box(slide, l, t, 0.07, h, TEAL)          # left accent bar
    box(slide, l, t, w, 0.42, NAVY)
    textbox(slide, title, l + 0.2, t + 0.05, w - 0.3, 0.35,
            size=title_size, bold=True, color=WHITE)
    tb = slide.shapes.add_textbox(
        Inches(l + 0.2), Inches(t + 0.52), Inches(w - 0.35), Inches(h - 0.62))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = DARK_GRAY


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide, NAVY)

# gold accent strip
box(slide, 0, 0, 0.18, 7.5, TEAL)
box(slide, 0, 5.6, 13.33, 0.06, GOLD)

textbox(slide, "Individual Achievement",
        0.5, 0.6, 12.3, 1.0, size=42, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
textbox(slide, "Presentation",
        0.5, 1.5, 12.3, 0.85, size=42, bold=True,
        color=GOLD, align=PP_ALIGN.LEFT)

box(slide, 0.5, 2.5, 6.0, 0.06, TEAL)

textbox(slide, "Tomin Jose",
        0.5, 2.7, 8.0, 0.65, size=26, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
textbox(slide, "Backend Lead  —  Project Orion",
        0.5, 3.35, 8.0, 0.5, size=18,
        color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)

textbox(slide, "AFL Player Tracking & Crowd Monitoring",
        0.5, 4.05, 10.0, 0.5, size=15,
        color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)
textbox(slide, "Expert Panel  ·  Week 6  ·  SIT782  ·  Trimester 1, 2026  ·  Deakin University",
        0.5, 5.75, 12.3, 0.45, size=12,
        color=MID_GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MY ROLE & OWNERSHIP
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide)
header_bar(slide, "My Role & Ownership",
           "Backend Lead — sole owner of the API Gateway from design to delivery")

# Left: role summary
box(slide, 0.3, 1.55, 5.9, 5.7, WHITE)
box(slide, 0.3, 1.55, 5.9, 0.42, DARK_RED)
textbox(slide, "What I owned", 0.45, 1.6, 5.6, 0.35,
        size=14, bold=True, color=WHITE)

role_items = [
    "▸  Designed the full microservices architecture",
    "▸  Built the FastAPI API Gateway end-to-end",
    "▸  Implemented auth, RBAC, job pipeline, upload",
    "▸  Orchestrated parallel async microservice calls",
    "▸  Reviewed & merged all team PRs",
    "▸  Authored BACKEND_PLAN.txt (architecture doc)",
    "▸  Created GIT_GUIDE.txt for team onboarding",
    "▸  Defined API contracts for frontend team",
]
tb = slide.shapes.add_textbox(Inches(0.45), Inches(2.1), Inches(5.6), Inches(4.8))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
for i, item in enumerate(role_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY

# Right: architecture diagram (simplified)
box(slide, 6.5, 1.55, 6.55, 5.7, WHITE)
box(slide, 6.5, 1.55, 6.55, 0.42, DARK_RED)
textbox(slide, "System I Designed", 6.65, 1.6, 6.2, 0.35,
        size=14, bold=True, color=WHITE)

box(slide, 7.8, 2.15, 3.8, 0.65, MID_NAVY)
textbox(slide, "Frontend  (React : 3000)", 7.8, 2.15, 3.8, 0.65,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(slide, "▼", 9.35, 2.82, 0.7, 0.4, size=20, color=TEAL, align=PP_ALIGN.CENTER)

box(slide, 7.5, 3.25, 4.4, 0.9, NAVY)
box(slide, 7.5, 3.25, 4.4, 0.07, GOLD)
textbox(slide, "API Gateway  (FastAPI : 8000)", 7.5, 3.32, 4.4, 0.45,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, "← MY RESPONSIBILITY →", 7.5, 3.72, 4.4, 0.35,
        size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

textbox(slide, "▼               ▼", 8.6, 4.2, 2.3, 0.4, size=18, color=TEAL)

box(slide, 6.8, 4.7, 2.6, 0.7, TEAL)
textbox(slide, "Player Service\n:8001", 6.8, 4.7, 2.6, 0.7,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(slide, 10.0, 4.7, 2.6, 0.7, MID_NAVY)
textbox(slide, "Crowd Service\n:8002", 10.0, 4.7, 2.6, 0.7,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(slide, 8.3, 5.6, 2.8, 0.65, RGBColor(0x2C, 0x3E, 0x50))
textbox(slide, "PostgreSQL", 8.3, 5.6, 2.8, 0.65,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — KEY ACHIEVEMENTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide)
header_bar(slide, "Key Technical Achievements",
           "What I built — and why it matters")

card(slide, "1.  Authentication & RBAC",
     ["✓  bcrypt password hashing — never stored in plaintext",
      "✓  JWT tokens — stateless, microservice-compatible",
      "✓  Role enforcement via FastAPI dependency injection",
      "✓  Two roles: user (own jobs) / admin (all jobs)"],
     0.3, 1.55, 6.2, 2.55)

card(slide, "2.  Video Upload & Job Pipeline",
     ["✓  Format validation — .mp4 / .avi / .mov only",
      "✓  BackgroundTasks — non-blocking video processing",
      "✓  4 job statuses: processing / done / partial / failed",
      "✓  Partial status — one service fails, results still saved"],
     6.7, 1.55, 6.3, 2.55)

card(slide, "3.  Parallel Microservice Orchestration",
     ["✓  asyncio.gather() — player + crowd called concurrently",
      "✓  Halves response latency vs sequential calls",
      "✓  Mock client pattern — decoupled from other teams",
      "✓  Retry endpoint — re-runs only the failed service"],
     0.3, 4.25, 6.2, 2.75)

card(slide, "4.  Code Review & Integration Leadership",
     ["✓  Reviewed & merged PR #15 (William — mock services)",
      "✓  Reviewed & merged PR #16 (Prabhnoor — logging)",
      "✓  Identified + reverted a conflicting logging impl",
      "✓  Maintained stable main branch throughout sprint"],
     6.7, 4.25, 6.3, 2.75)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GOING ABOVE & BEYOND (HD DIFFERENTIATORS)
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide)
header_bar(slide, "Going Above & Beyond",
           "Initiatives that exceed the standard contribution — HD differentiators")

items = [
    ("BACKEND_PLAN.txt — Architecture Decision Log",
     "Documented every major technical decision with rationale before writing a single "
     "line of code. Covered DB choice, auth strategy, async pattern, job status design, "
     "and phased delivery. Served as the team's single source of truth and prevented "
     "repeated direction discussions throughout the sprint.",
     0.3, 1.55),
    ("Mock Client Pattern — Unblocking the Team",
     "Introduced stub player_client.py and crowd_client.py early in development so the "
     "full backend pipeline was testable weeks before real microservice contracts were "
     "finalised. Inverted the dependency: I defined the interface and let other teams "
     "implement to it — a professional-grade decoupling strategy.",
     0.3, 3.3),
    ("Partial Job Status — Fault-Tolerant Design",
     "Recognised that treating any downstream failure as a full job failure would lose "
     "valid results. Designed the 'partial' status to preserve and return successful "
     "service output even when the other fails. This required reasoning about system "
     "reliability from the user's perspective, not just implementation convenience.",
     0.3, 5.05),
]

for title, body, x, y in items:
    box(slide, x, y, 12.7, 1.6, WHITE)
    box(slide, x, y, 0.07, 1.6, TEAL)
    textbox(slide, title, x + 0.35, y + 0.08, 12.1, 0.38,
            size=14, bold=True, color=DARK_RED)
    textbox(slide, body, x + 0.35, y + 0.52, 12.1, 1.0,
            size=12, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — IMPACT ON THE PROJECT
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide)
header_bar(slide, "Impact on the Project",
           "How my contributions enabled the rest of the team")

impacts = [
    ("Frontend Team",
     ["Stable API contracts defined early",
      "OpenAPI docs at /docs — self-service reference",
      "No blocking dependencies on my delivery"],
     TEAL, 0.3, 1.55),
    ("Player & Crowd Teams",
     ["Service interface template provided",
      "Mock clients mean integration is drop-in",
      "Gateway already handles partial failure"],
     MID_NAVY, 4.65, 1.55),
    ("Backend Team Members",
     ["BACKEND_PLAN.txt — onboarding without me",
      "GIT_GUIDE.txt — consistent branch workflow",
      "PR reviews raised team code quality"],
     GOLD, 8.98, 1.55),
]

for title, bullets, color, x, y in impacts:
    box(slide, x, y, 4.1, 3.5, LIGHT_BLUE)
    box(slide, x, y, 4.1, 0.42, color)
    textbox(slide, title, x + 0.15, y + 0.05, 3.8, 0.35,
            size=14, bold=True, color=WHITE if color != GOLD else NAVY)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.55), Inches(3.8), Inches(2.7))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

# Summary bar
box(slide, 0.3, 5.3, 12.7, 1.9, NAVY)
box(slide, 0.3, 5.3, 0.07, 1.9, TEAL)
textbox(slide, "Overall Impact",
        0.5, 5.35, 3.0, 0.45, size=15, bold=True, color=GOLD)
textbox(slide,
        "The backend gateway is the central integration point for all four sub-teams. "
        "By delivering a stable, documented, and extensible foundation ahead of schedule, "
        "I removed the #1 risk to project delivery — inter-team dependency delays — "
        "and gave every other team a concrete surface to build and test against.",
        0.5, 5.82, 12.3, 1.3, size=13, color=RGBColor(0xAA, 0xCC, 0xEE))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — REFLECTION & LEARNING
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide)
header_bar(slide, "Reflection & Learning",
           "What this experience has taught me")

reflections = [
    ("Architecture First Pays Off",
     "Spending time on BACKEND_PLAN.txt before coding prevented significant "
     "rework. Every 'why did we do it this way?' question during the sprint "
     "had a documented answer. I'll apply this discipline to every future "
     "project — especially in team settings where alignment is expensive.",
     "GLO1 / SWDN"),
    ("Leading Through Code Review",
     "Managing the revert of Prabhnoor's logging PR taught me that good "
     "leadership means making technically correct decisions even when it "
     "creates short-term friction. The result — a clean, stable codebase — "
     "validated that approach.",
     "GLO7"),
    ("Designing for Failure, Not Just Success",
     "The 'partial' job status was a conscious design choice driven by thinking "
     "about what users actually experience in degraded conditions. This shift — "
     "from 'does it work?' to 'how does it fail?' — is the most important "
     "engineering mindset I've developed this trimester.",
     "GLO4 / GLO5"),
    ("Decoupling as a Team Strategy",
     "The mock client pattern was as much a team coordination tool as a "
     "technical one. Defining interfaces early and letting other teams implement "
     "to them removed a critical path dependency. I now see interface design "
     "as a leadership decision, not just a technical one.",
     "GLO2 / GLO7"),
]

for i, (title, body, glo) in enumerate(reflections):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.55 + row * 2.85
    box(slide, x, y, 6.3, 2.65, WHITE)
    box(slide, x, y, 6.3, 0.42, DARK_RED)
    textbox(slide, title, x + 0.15, y + 0.05, 5.5, 0.35,
            size=13, bold=True, color=WHITE)
    box(slide, x + 5.3, y + 0.05, 0.85, 0.32, TEAL)
    textbox(slide, glo, x + 5.3, y + 0.05, 0.85, 0.32,
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, body, x + 0.15, y + 0.52, 6.0, 2.0,
            size=12, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT'S NEXT
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide)
header_bar(slide, "What's Next",
           "Remaining trimester — completing the vision")

next_items = [
    ("Real Microservice Integration",
     ["Replace mock clients with live player/crowd endpoints",
      "Formal contract review before each integration",
      "Validate response schemas against gateway expectations"],
     0.3, 1.55, 4.1),
    ("End-to-End Integration Testing",
     ["Full upload → process → result pipeline tests",
      "Test partial failure and retry behaviour",
      "Cover edge cases: bad video, service timeout"],
     4.65, 1.55, 4.1),
    ("Docker Compose Deployment",
     ["Containerise all 4 services for demo submission",
      "Define networking and volume configuration",
      "Co-ordinate with each sub-team on containers"],
     8.98, 1.55, 4.1),
]

for title, bullets, x, y, w in next_items:
    box(slide, x, y, w, 3.3, WHITE)
    box(slide, x, y, w, 0.42, DARK_RED)
    textbox(slide, title, x + 0.15, y + 0.05, w - 0.3, 0.35,
            size=13, bold=True, color=WHITE)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.55), Inches(w - 0.3), Inches(2.5))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

# Bottom banner
box(slide, 0.3, 5.1, 12.7, 2.1, NAVY)
box(slide, 0.3, 5.1, 0.07, 2.1, GOLD)
textbox(slide, "HD Commitment",
        0.5, 5.15, 3.0, 0.45, size=15, bold=True, color=GOLD)
textbox(slide,
        "I will continue to lead backend delivery with the same standard of "
        "technical rigour, documented decision-making, and proactive team "
        "enablement that has characterised my work so far. My goal is not just "
        "a working system — but a system the whole team is proud to demonstrate.",
        0.5, 5.65, 12.3, 1.4, size=13, color=RGBColor(0xAA, 0xCC, 0xEE))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
bg(slide, NAVY)
box(slide, 0, 0, 0.18, 7.5, TEAL)
box(slide, 0, 2.95, 13.33, 0.06, GOLD)

textbox(slide, "Thank You",
        0.5, 0.35, 12.3, 1.3, size=56, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
textbox(slide, "Tomin Jose  —  Backend Lead, Project Orion",
        0.5, 1.65, 12.3, 0.55, size=18,
        color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)
textbox(slide, "Questions & Feedback Welcome",
        0.5, 2.25, 12.3, 0.45, size=15,
        color=GOLD, align=PP_ALIGN.LEFT)

# Left summary card
box(slide, 0.3, 3.15, 6.1, 4.0, RGBColor(0x12, 0x3A, 0x70))
box(slide, 0.3, 3.15, 6.1, 0.42, TEAL)
textbox(slide, "Achievements Summary",
        0.45, 3.2, 5.8, 0.35, size=13, bold=True, color=WHITE)
summary = (
    "▸  Full API Gateway — auth, jobs, upload, orchestration\n"
    "▸  RBAC + JWT security system\n"
    "▸  Parallel async microservice calls\n"
    "▸  Partial failure tolerance (partial job status)\n"
    "▸  Mock client pattern — unblocked 3 sub-teams\n"
    "▸  PR reviews + stable main branch leadership"
)
textbox(slide, summary, 0.45, 3.65, 5.8, 3.3, size=13,
        color=RGBColor(0xCC, 0xDD, 0xEE))

# Right GLO card
box(slide, 6.8, 3.15, 6.2, 4.0, RGBColor(0x12, 0x3A, 0x70))
box(slide, 6.8, 3.15, 6.2, 0.42, GOLD)
textbox(slide, "GLOs Demonstrated",
        6.95, 3.2, 5.9, 0.35, size=13, bold=True, color=NAVY)
glos = (
    "▸  GLO1 — Expert FastAPI / backend knowledge\n"
    "▸  GLO2 — Cross-team API contract communication\n"
    "▸  GLO4 — Critical design decisions (partial status)\n"
    "▸  GLO5 — Mock pattern for ill-defined problem\n"
    "▸  GLO6 — Self-managed phased delivery\n"
    "▸  GLO7 — Led team through PR reviews"
)
textbox(slide, glos, 6.95, 3.65, 5.9, 3.3, size=13,
        color=RGBColor(0xCC, 0xDD, 0xEE))


# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = (
    "d:/Masters/trimester 1 2026/SIT782 \u2013 Team Project (B) \u2013 Execution and Delivery"
    "/redback-orion/26_T1/afl_player_tracking_and_crowd_monitoring/backend/5.4hd.pptx"
)
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
